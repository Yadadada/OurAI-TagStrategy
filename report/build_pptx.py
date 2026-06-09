"""
生成原生可编辑的 .pptx —— 每个文字框、表格都可以在 PowerPoint 里点开改字。
不依赖 Marp。
用法：python build_pptx.py
输出：report_editable.pptx（同目录）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ---------- 颜色配置 ----------
NAVY = RGBColor(0x1A, 0x36, 0x5D)
BLUE = RGBColor(0x2C, 0x55, 0x82)
RED = RGBColor(0xC5, 0x30, 0x30)
LIGHT_BLUE = RGBColor(0xEB, 0xF4, 0xFF)
GRAY_BG = RGBColor(0xF7, 0xFA, 0xFC)
DARK = RGBColor(0x2D, 0x37, 0x48)

# ---------- 创建 16:9 ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="微软雅黑"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_title(slide, text, *, size=30, color=NAVY, y=Inches(0.4)):
    add_text(slide, Inches(0.6), y, Inches(12), Inches(0.7), text, size=size, bold=True, color=color)
    # underline bar
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y + Inches(0.62), Inches(12), Emu(20000))
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY


def add_footer(slide, page_num, total=8):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             "OurAI 问卷自由文本标签抽取 · 2026-06", size=10, color=RGBColor(0x80, 0x80, 0x80))
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
             f"{page_num} / {total}", size=10, color=RGBColor(0x80, 0x80, 0x80), align=PP_ALIGN.RIGHT)


def add_table(slide, x, y, w, h, data, *, header=True, col_widths=None, font_size=14, header_bg=LIGHT_BLUE):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = table_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = int(w * cw / total)
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.name = "微软雅黑"
            run.font.size = Pt(font_size)
            if header and r == 0:
                run.font.bold = True
                run.font.color.rgb = NAVY
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            else:
                run.font.color.rgb = DARK
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return table


def add_bullet_box(slide, x, y, w, h, items, *, size=16, color=DARK):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "• " + item
        run.font.name = "微软雅黑"
        run.font.size = Pt(size)
        run.font.color.rgb = color


def add_highlight_box(slide, x, y, w, h, text, *, bg=RGBColor(0xFF, 0xF5, 0xCC), size=15):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = RGBColor(0xE2, 0xB7, 0x00)
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "微软雅黑"
    run.font.size = Pt(size)
    run.font.color.rgb = DARK


def set_speaker_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


# ============ P1 封面 ============
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()

add_text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(1.2),
         "OurAI 问卷自由文本标签抽取", size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(3.1), Inches(12), Inches(0.7),
         "固定主标签 + 动态子标签  ·  Prompt 调优  ·  评测驱动", size=22, color=RGBColor(0xCB, 0xD5, 0xE0), align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(4.3), Inches(12), Inches(0.5),
         "从关键词规则 53.3% 到 LLM 零样本 97.5%  ·  同分布偏差双重验证",
         size=16, color=RGBColor(0xA0, 0xAE, 0xC0), align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(6.0), Inches(12), Inches(0.5),
         "汇报人:项目组   ·   2026-06", size=14, color=RGBColor(0xA0, 0xAE, 0xC0), align=PP_ALIGN.CENTER)
set_speaker_notes(s, "开场：今天汇报 OurAI 问卷自由文本标签抽取子模块。重点三件事：为什么要引入这套方案、Prompt 怎么调优、怎么评测且防偏差。5 分钟，欢迎随时打断。")

# ============ P2 问题背景 ============
s = prs.slides.add_slide(BLANK)
add_title(s, "一 · 问题背景")

add_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.4),
         "OurAI 问卷有三个开放性文本字段，天然不可计算", size=16, bold=True, color=DARK)
add_table(s, Inches(0.6), Inches(1.85), Inches(12), Inches(1.85),
          [
              ["字段", "问题", "下游用途"],
              ["Q22  intro_prompt", "第一次见面，你最希望一起做什么？", "约会场景匹配"],
              ["Q23  q19", "你受不了对方什么？", "关系雷区匹配"],
              ["Q24  q20", "对 TA 有什么补充要求？", "择偶偏好匹配"],
          ],
          col_widths=[2.5, 5, 3], font_size=14)

add_text(s, Inches(0.6), Inches(3.9), Inches(12), Inches(0.4),
         "两种极端方案都走不通", size=16, bold=True, color=DARK)
add_table(s, Inches(0.6), Inches(4.4), Inches(12), Inches(1.5),
          [
              ["方案", "问题"],
              ["纯关键词规则", "fallback 实测主标签准确率仅 53.3%，网络词/否定句/并列结构全漏"],
              ["纯开放文本直接入库", "无法向量化检索，下游匹配推荐无法消费"],
          ],
          col_widths=[3, 8], font_size=14)

add_highlight_box(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.8),
                  "选型：固定主标签（保证可计算）+ 动态子标签（保留语义细节）+ LLM 抽取 + 关键词降级兜底",
                  size=15)
add_footer(s, 2)
set_speaker_notes(s, "53.3% 是真实跑过 fallback 的数字。中央空调、探店、爹味这种网络词，规则根本抓不到。纯开放文本又没法给匹配算法消费。所以引出这套方案。")

# ============ P3 技术方案 ============
s = prs.slides.add_slide(BLANK)
add_title(s, "二 · 技术方案")

add_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.4),
         "标签体系：三字段各自的主标签设计", size=16, bold=True, color=DARK)
add_table(s, Inches(0.6), Inches(1.85), Inches(12), Inches(2.0),
          [
              ["字段", "主标签", "设计来源"],
              ["Q22 场景", "5 类 + 互动横切维度（强 / 弱 / null）", "Z 世代约会消费数据 + 校园场景聚类"],
              ["Q23 雷区", "5 类", "Gottman「四骑士」+ 中文网络词聚类"],
              ["Q24 补充", "7 类", "Buss 择偶偏好量表 + 中文社交平台真实择偶帖"],
          ],
          col_widths=[2, 4.5, 5], font_size=14)

add_text(s, Inches(0.6), Inches(4.05), Inches(12), Inches(0.4),
         "抽取流程", size=16, bold=True, color=DARK)

# 流程图（inline，在技术方案页）
def flow_box_s(slide, x, y, w, h, text, color=BLUE, font_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = NAVY
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "微软雅黑"
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

def arrow_s(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)
    conn.line.color.rgb = NAVY
    conn.line.width = Pt(1.2)

y_flow = Inches(4.6)
flow_items = [
    (Inches(0.5),  "输入文本", BLUE),
    (Inches(2.2),  "Prompt\n构造", BLUE),
    (Inches(4.2),  "LLM 调用\nOpenAI 兼容", BLUE),
    (Inches(6.5),  "主标签\n白名单校验", BLUE),
    (Inches(8.8),  "输出\n结构化标签", RGBColor(0x38, 0xA1, 0x69)),
]
for x, txt, col in flow_items:
    flow_box_s(s, x, y_flow, Inches(1.6), Inches(0.85), txt, col)
for i in range(len(flow_items) - 1):
    x1 = flow_items[i][0] + Inches(1.6)
    x2 = flow_items[i+1][0]
    arrow_s(s, x1, y_flow + Inches(0.42), x2, y_flow + Inches(0.42))

# Fallback 分支标注
flow_box_s(s, Inches(4.2), Inches(5.7), Inches(2.0), Inches(0.65), "↓ API 失败\n关键词降级兜底", RGBColor(0x80, 0x80, 0x80), font_size=10)

add_text(s, Inches(11.0), Inches(4.6), Inches(2.2), Inches(1.8),
         "三个运行时开关\nFEW_SHOT=1\nCOT=1\n自动降级",
         size=11, color=BLUE)
add_footer(s, 3)
set_speaker_notes(s, "互动横切维度是 Q22 特有的，不替代主标签，叠加判断双人对打/安静陪伴/无信号。白名单校验保证 LLM 乱输也不会污染下游。三个开关都是环境变量，不需要重新部署。每个主标签下挂 2-3 个预设子标签，LLM 可以输出词表外的子标签，用 quote 字段锚回原文。")

# ============ P4 Prompt 调优实验设计 ============
s = prs.slides.add_slide(BLANK)
add_title(s, "三 · Prompt 调优 — 实验设计")

add_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.4),
         "攻关目标：Q23 longtail 桶基线 76.5%（网络词：中央空调 / PUA / 时间管理大师——分类混乱）",
         size=15, bold=True, color=RED)

add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.4),
         "三组实验（同一 244 条评测集，剔除 few-shot 样本保证公平）", size=15, color=DARK)
add_table(s, Inches(0.6), Inches(2.5), Inches(12), Inches(2.7),
          [
              ["实验", "配置", "核心思路"],
              ["#4  零样本", "当前 prompt，无额外示例", "锚点基线"],
              ["#2  +few-shot", "每字段 4 条示例（normal / longtail / boundary / empty 四象限）", "教格式 + 给网络词范例"],
              ["#3  +few-shot+CoT", "在 #2 基础上加 reasoning 字段，先翻译词义再选标签", "让模型显式推理\n中央空调→不专一→dishonesty"],
          ],
          col_widths=[2.5, 5.5, 4], font_size=13)

add_highlight_box(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.75),
                  "公平性保障：12 条 few-shot 样本硬编码 ID，评测时用 --exclude-few-shot 标志剔除（256→244）",
                  size=14)

add_text(s, Inches(0.6), Inches(6.45), Inches(12), Inches(0.4),
         "few-shot 真正的价值：子标签命中率 Q22 +19pp（教会模型怎么从原文提炼 sub 字段）",
         size=13, color=BLUE)
add_footer(s, 4)
set_speaker_notes(s, "公平性设计很重要：12 条 few-shot 样本硬编码 ID 从评测集剔除，不然 #2 的准确率会虚高。CoT 的核心假设是让模型先说出语义链再分类，对词义清晰的网络词应该有效。few-shot 真正的价值不在主标签，而在子标签。")

# ============ P5 Prompt 调优结果 ============
s = prs.slides.add_slide(BLANK)
add_title(s, "三 · Prompt 调优 — 三轴结果与反直觉发现")

add_table(s, Inches(0.6), Inches(1.35), Inches(12), Inches(2.2),
          [
              ["实验", "主标签准确率", "平均耗时", "Token / 样本", "成本相对基线"],
              ["#4  零样本", "97.5%", "1464ms", "732", "基准"],
              ["#2  +few-shot", "96.7%", "1603ms", "1036", "+41%"],
              ["#3  +few-shot+CoT", "97.5%", "1573ms", "1154", "+58%"],
          ],
          col_widths=[3, 2.5, 2, 2, 2.5], font_size=14)

add_highlight_box(s, Inches(0.6), Inches(3.8), Inches(12), Inches(0.65),
                  "⚠  反直觉发现：零样本最便宜，且与 #3 并列最准",
                  bg=RGBColor(0xFE, 0xF3, 0xC7), size=16)

add_text(s, Inches(0.6), Inches(4.65), Inches(12), Inches(0.35),
         "为什么 few-shot / CoT 没赢：", size=14, bold=True, color=DARK)
add_bullet_box(s, Inches(0.8), Inches(5.05), Inches(11.5), Inches(1.3), [
    "few-shot 把 304 tok 固定堆在每次 prompt，主标签准确率没有对应提升",
    "CoT 仅修对 2 条 longtail，同时让互动维度准确率倒退 8.7pp（注意力分散副作用）",
], size=13)

add_table(s, Inches(0.6), Inches(6.35), Inches(12), Inches(0.9),
          [
              ["极致省钱 → #4 零样本", "生产平衡 → #2 few-shot", "longtail 高发字段 → 按字段开 CoT"],
          ],
          col_widths=[1, 1, 1], font_size=13, header=False)
add_footer(s, 5)
set_speaker_notes(s, "加成本这第三轴才看出真相——原来只看精度+速度，差点认定 #2 是 winner。GPT-4.1 本身已经吸收了大量分类任务训练数据，多塞 4 条 few-shot 是在教它已经会的事。调优前先跑零样本基线，别一上来就 few-shot+CoT，这是最值得带走的方法论。")

# ============ P6 评测数据集构造 ============
s = prs.slides.add_slide(BLANK)
add_title(s, "四 · 评测体系 — 数据集构造")

add_text(s, Inches(0.6), Inches(1.35), Inches(12), Inches(0.4),
         "三段式方法（总计 256 条，q22=85 / q23=87 / q24=84）", size=15, bold=True, color=DARK)

# 三段流程图
y_f = Inches(1.9)
flow_box_s(s, Inches(0.6),  y_f, Inches(2.4), Inches(0.85), "人工种子\n46 条", BLUE, font_size=13)
arrow_s(s, Inches(3.0), y_f + Inches(0.42), Inches(3.8), y_f + Inches(0.42))
flow_box_s(s, Inches(3.8),  y_f, Inches(3.2), Inches(0.85), "GPT-4.1 网格扩充\n210 条", BLUE, font_size=13)
arrow_s(s, Inches(7.0), y_f + Inches(0.42), Inches(7.8), y_f + Inches(0.42))
flow_box_s(s, Inches(7.8),  y_f, Inches(2.4), Inches(0.85), "10% 人工抽检\n22 条，合格率 100%", BLUE, font_size=13)
arrow_s(s, Inches(10.2), y_f + Inches(0.42), Inches(11.0), y_f + Inches(0.42))
flow_box_s(s, Inches(11.0), y_f, Inches(1.8), Inches(0.85), "256 条\n样本池", RGBColor(0x38, 0xA1, 0x69), font_size=13)

add_text(s, Inches(0.6), Inches(3.0), Inches(12), Inches(0.4),
         "关键设计：按 (主标签 × category) 网格调度，强制扩充比例", size=14, bold=True, color=DARK)
add_table(s, Inches(0.6), Inches(3.5), Inches(12), Inches(2.6),
          [
              ["category", "比例", "目的"],
              ["normal  普通典型", "60%", "让整体准确率有 baseline 意义"],
              ["longtail  网络词", "20%", "Prompt 调优有足够攻关样本"],
              ["boundary  双主标签", "10%", "测试边界判断能力"],
              ["perturbation  否定/反问", "10%", "测试语义鲁棒性"],
              ["empty  敷衍/不答", "种子覆盖", "LLM 生成的空回答不可信"],
          ],
          col_widths=[3, 1.5, 7], font_size=13)
add_highlight_box(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.65),
                  "不做网格调度时 LLM 偷懒只生成 normal，longtail 样本严重不足，Prompt 调优无从下手",
                  size=13)
add_footer(s, 6)
set_speaker_notes(s, "为什么要网格调度：自由生成时 GPT-4.1 会全给 normal，longtail 根本没有足够样本量，prompt 调优就无从下手。empty 类型有限不靠 LLM 生成是个坑：LLM 生成的空回答往往是随便写一个短句，不是真实用户的敷衍模式。")

# ============ P7 评测指标 + 同分布偏差双重验证 ============
s = prs.slides.add_slide(BLANK)
add_title(s, "四 · 评测体系 — 指标 + 同分布偏差双重验证")

# 左栏：指标
add_text(s, Inches(0.6), Inches(1.35), Inches(5.8), Inches(0.4),
         "4 个核心指标", size=14, bold=True, color=DARK)
add_table(s, Inches(0.6), Inches(1.8), Inches(5.8), Inches(2.0),
          [
              ["指标", "说明"],
              ["主标签准确率", "∈ expected_main 即对"],
              ["子标签宽松命中率", "substring 双向匹配"],
              ["互动维度准确率", "仅 Q22，三值严格相等"],
              ["空判召回率", "should_skip 时必须为空"],
          ],
          col_widths=[3, 3], font_size=12)

# 右栏：同分布偏差
add_text(s, Inches(7.0), Inches(1.35), Inches(6.1), Inches(0.4),
         "同分布偏差双重验证", size=14, bold=True, color=RED)

add_text(s, Inches(7.0), Inches(1.85), Inches(6.1), Inches(0.35),
         "第一层：跨模型对照（同 256 条）", size=13, bold=True, color=DARK)
add_table(s, Inches(7.0), Inches(2.25), Inches(6.1), Inches(1.5),
          [
              ["模型", "准确率", "差距"],
              ["GPT-4.1", "97.5%", "—"],
              ["Claude 4.5/4.7", "~95-97%", "≤ 3pp"],
              ["Gemini", "~95-97%", "≤ 3pp"],
          ],
          col_widths=[3, 2, 2], font_size=12)

add_text(s, Inches(7.0), Inches(3.9), Inches(6.1), Inches(0.35),
         "第二层：held-out 对照集（80 条异源 + longtail 45%）", size=13, bold=True, color=DARK)
add_table(s, Inches(7.0), Inches(4.3), Inches(6.1), Inches(1.3),
          [
              ["数据集", "主标签准确率"],
              ["扩充集 256 条", "97.5%"],
              ["Held-out 80 条", "92.5%"],
          ],
          col_widths=[4, 2], font_size=12)
add_text(s, Inches(7.0), Inches(5.75), Inches(6.1), Inches(0.4),
         "6 条错全部是 longtail 标注分歧，非系统失败\nnormal / boundary / perturbation 全部 100%",
         size=12, color=DARK)

add_highlight_box(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.65),
                  "三家模型差距 ≤ 3pp + held-out 92.5%  →  同分布偏差在当前规模下可控，97.5% 不是自洽虚高",
                  bg=RGBColor(0xC6, 0xF6, 0xD5), size=13)
add_footer(s, 7)
set_speaker_notes(s, "负责人直接质疑同分布偏差——这是个真问题。如果只用 GPT-4.1 评测自己生成的数据，97.5% 可能是虚高。第一层：换三家模型跑，差距 ≤ 3pp，说明评测集通用。第二层：换生成源(Claude)、100% 人工筛选、longtail 比例翻倍，三重压力下还有 92.5%，且失败模式可解释。")

# ============ P8 结论 ============
s = prs.slides.add_slide(BLANK)
add_title(s, "五 · 结论")

add_table(s, Inches(0.6), Inches(1.35), Inches(12), Inches(2.5),
          [
              ["维度", "结果"],
              ["方案可行性", "关键词规则 53.3%  →  LLM 零样本 97.5%，提升 44pp"],
              ["Prompt 调优", "零样本已是 Pareto 最优；few-shot 的真实价值在子标签（+19pp），不在主标签"],
              ["评测可信度", "双重验证确认同分布偏差可控，held-out 92.5% 是保守下界"],
          ],
          col_widths=[2.5, 9.5], font_size=15)

add_text(s, Inches(0.6), Inches(4.1), Inches(12), Inches(0.4),
         "最值得复用的工程资产", size=16, bold=True, color=DARK)
add_highlight_box(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.75),
                  "三段式评测数据集（种子→网格扩充→抽检）+ 5 类分桶指标 + ADR 决策记录  →  可迁移到任意标签抽取任务",
                  size=14)

add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.4),
         "未来工作", size=16, bold=True, color=DARK)
add_bullet_box(s, Inches(0.8), Inches(6.1), Inches(11.5), Inches(1.0), [
    "跨领域迁移（DomainConfig 注入，验证电商 / 游戏场景）",
    "工程化：请求队列 + 熔断 + Token usage 可观测性  ·  应用：匹配推荐赋能 / 冷启动对话",
], size=13)
add_footer(s, 8)
set_speaker_notes(s, "结束语：单一的 97.5% 没什么了不起，关键是这套评测驱动+决策可追溯+偏差防护的方法本身可以迁到任何标签抽取任务上。Q&A 时间。")

# ---------- 保存 ----------
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "report_editable.pptx")
prs.save(out)
print(f"[OK] generated: {out}")
print(f"     {len(prs.slides)} slides, all editable in PowerPoint")
